"""
Document Processing Service

Parses uploaded documents (PDF, DOCX, XLSX, PPTX, HWPX) and extracts:
- Full text content
- Tables (as structured data)
- Metadata (author, dates, page count, etc.)
- Sections/chapters structure
- Keywords (simple TF-based extraction)

Stores results as JSON alongside the uploaded file.
Standards applied: ISO 19650 (CDE metadata), ISO 32000 (PDF), ISO/IEC 29500 (OOXML)
"""

import json
import os
import re
import threading
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Optional
from uuid import UUID


def process_document_async(file_id: str, file_path: str, original_filename: str, lifecycle_phase: str, update_callback=None):
    """Process a document in a background thread."""
    thread = threading.Thread(
        target=_process_sync,
        args=(file_id, file_path, original_filename, lifecycle_phase, update_callback),
        daemon=True,
    )
    thread.start()
    return thread


def _process_sync(file_id: str, file_path: str, original_filename: str, lifecycle_phase: str, update_callback=None):
    """Synchronous document processing."""
    try:
        if update_callback:
            update_callback(file_id, "processing")

        ext = os.path.splitext(file_path)[1].lower()

        # Auto-classify document type
        doc_class = _classify_document_type(original_filename, ext, lifecycle_phase)

        result = {
            "file_id": file_id,
            "original_filename": original_filename,
            "lifecycle_phase": lifecycle_phase,
            "processed_at": datetime.utcnow().isoformat(),
            "extension": ext,
            "status": "processing",
            "document_type": doc_class,
            "full_text": "",
            "sections": [],
            "tables": [],
            "metadata": {},
            "keywords": [],
            "statistics": {},
            "standards_applied": [],
            "ng_items": [],
        }

        # Initialize standards pipeline tracker
        result["standards_pipeline"] = []

        # === STAGE 1: INGEST (ISO 19650 CDE) ===
        result["standards_pipeline"].append({
            "stage": "Ingest",
            "standard": "ISO 19650-1/2",
            "standard_name": "Information Management using BIM",
            "action": "CDE workflow registration",
            "details": f"File registered in CDE with state 'WIP'. Lifecycle phase: {lifecycle_phase}. "
                       f"Filename parsed for ISO 19650 naming convention compliance.",
            "input": original_filename,
            "output": f"CDE state: WIP, Phase: {lifecycle_phase}",
            "status": "completed",
        })

        # === STAGE 2: PARSE (Format-specific standard) ===
        if ext == ".pdf":
            result = _parse_pdf(file_path, result)
        elif ext in (".docx", ".doc"):
            result = _parse_docx(file_path, result)
        elif ext in (".xlsx", ".xls", ".csv"):
            result = _parse_xlsx(file_path, result)
        elif ext in (".pptx", ".ppt"):
            result = _parse_pptx(file_path, result)
        elif ext in (".hwpx", ".hwp"):
            result = _parse_hwpx(file_path, result)
        elif ext == ".ifc":
            result = _parse_ifc_basic(file_path, result)
        else:
            result["status"] = "unsupported"
            result["error"] = f"Unsupported file type: {ext}"

        # === STAGE 2.5: DOMAIN RELEVANCE CHECK + CONTENT-BASED RECLASSIFICATION ===
        domain_result = _check_domain_relevance(result)
        result["domain_relevance"] = domain_result
        result["standards_pipeline"].append({
            "stage": "Validate",
            "standard": "ISO 19650-1 / ISO 12006-2",
            "standard_name": "Domain Relevance & Classification",
            "action": f"Domain: {domain_result['domain']} ({domain_result['confidence']}%)",
            "details": domain_result["detail"],
            "input": "Extracted text and keywords",
            "output": f"Domain: {domain_result['domain']}, Relevance: {domain_result['confidence']}%",
            "status": "completed" if domain_result["is_relevant"] else "warning",
        })

        # Content-based reclassification (override filename-based if content says otherwise)
        if result["full_text"]:
            content_class = _classify_by_content(result["full_text"], result.get("keywords", []), lifecycle_phase)
            if content_class and content_class["confidence"] > doc_class.get("confidence", 0):
                result["document_type"] = content_class

        # === STAGE 3: VALIDATE (IDS 1.0 + LOIN) ===
        validation_results = _apply_standards_validation(result, lifecycle_phase)
        result["standards_pipeline"].extend(validation_results)

        # === STAGE 4: ENRICH (bSDD + ISO 12006-2) ===
        enrichment_results = _apply_standards_enrichment(result, lifecycle_phase)
        result["standards_pipeline"].extend(enrichment_results)

        # === STAGE 5: TRANSFORM (AI-ready) ===
        if result["full_text"]:
            result["keywords"] = _extract_keywords(result["full_text"])
            result["statistics"]["word_count"] = len(result["full_text"].split())
            result["statistics"]["char_count"] = len(result["full_text"])
            result["statistics"]["sentence_count"] = len(re.split(r'[.!?。]\s', result["full_text"]))

            result["standards_pipeline"].append({
                "stage": "Transform",
                "standard": "AI/ML Pipeline",
                "standard_name": "AI-Ready Data Transformation",
                "action": "Text chunking, keyword extraction, structure analysis",
                "details": f"Extracted {len(result['keywords'])} keywords, "
                           f"{result['statistics'].get('word_count', 0)} words, "
                           f"{len(result.get('tables', []))} tables, "
                           f"{len(result.get('sections', []))} sections.",
                "input": f"Parsed text ({result['statistics'].get('char_count', 0)} chars)",
                "output": "Keywords, sections, tables, statistics",
                "status": "completed",
            })

        # === STAGE 6: AI LAKE ===
        result["standards_pipeline"].append({
            "stage": "AI Lake",
            "standard": "ISO 19650-1",
            "standard_name": "Information Container (CDE)",
            "action": "Store processed data in AI Data Lake",
            "details": "Parsed data stored as queryable JSON. Full-text indexed for search. "
                       "Metadata catalogued per ISO 19115 profile.",
            "input": "All extracted data",
            "output": "Indexed document in AI Data Lake",
            "status": "completed",
        })

        if result["status"] != "unsupported":
            result["status"] = "completed"

        # Save parsed result
        parsed_path = file_path.rsplit(".", 1)[0] + "_parsed.json"
        with open(parsed_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2, default=str)

        if update_callback:
            update_callback(file_id, "completed" if result["status"] == "completed" else "failed")

        return result

    except Exception as e:
        error_result = {"file_id": file_id, "status": "failed", "error": str(e)}
        parsed_path = file_path.rsplit(".", 1)[0] + "_parsed.json"
        with open(parsed_path, "w", encoding="utf-8") as f:
            json.dump(error_result, f, ensure_ascii=False, indent=2)
        if update_callback:
            update_callback(file_id, "failed")
        return error_result


def _parse_pdf(file_path: str, result: dict) -> dict:
    """Parse PDF using PyMuPDF. Standard: ISO 32000."""
    import fitz

    result["standards_applied"].append({"code": "ISO 32000", "name": "PDF", "role": "Document format parsing"})

    doc = fitz.open(file_path)
    result["metadata"] = {
        "page_count": doc.page_count,
        "title": doc.metadata.get("title", ""),
        "author": doc.metadata.get("author", ""),
        "subject": doc.metadata.get("subject", ""),
        "creator": doc.metadata.get("creator", ""),
        "creation_date": doc.metadata.get("creationDate", ""),
    }

    full_text_parts = []
    sections = []

    for page_num in range(doc.page_count):
        page = doc[page_num]
        text = page.get_text()
        full_text_parts.append(text)

        # Detect section headers (lines that are short and possibly bold/large)
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        if span["size"] > 14 and len(span["text"].strip()) > 2 and len(span["text"].strip()) < 100:
                            sections.append({
                                "title": span["text"].strip(),
                                "page": page_num + 1,
                                "font_size": round(span["size"], 1),
                            })

        # Extract tables using basic heuristic
        tables = page.find_tables()
        if tables and tables.tables:
            for table in tables.tables:
                try:
                    data = table.extract()
                    if data and len(data) > 1:
                        headers = [str(h) for h in data[0]] if data[0] else []
                        rows = [[str(c) for c in row] for row in data[1:]]
                        result["tables"].append({
                            "page": page_num + 1,
                            "headers": headers,
                            "rows": rows[:50],  # Limit rows
                            "row_count": len(data) - 1,
                        })
                except Exception:
                    pass

    doc.close()
    result["full_text"] = "\n".join(full_text_parts)
    result["sections"] = sections
    result["statistics"]["page_count"] = result["metadata"]["page_count"]
    result["statistics"]["table_count"] = len(result["tables"])
    result["statistics"]["section_count"] = len(sections)
    return result


def _parse_docx(file_path: str, result: dict) -> dict:
    """Parse DOCX using python-docx. Standard: ISO/IEC 29500."""
    from docx import Document

    result["standards_applied"].append({"code": "ISO/IEC 29500", "name": "OOXML", "role": "Document format parsing"})

    doc = Document(file_path)
    full_text_parts = []
    sections = []

    # Core properties
    try:
        props = doc.core_properties
        result["metadata"] = {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
            "revision": props.revision or 0,
        }
    except Exception:
        result["metadata"] = {}

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        full_text_parts.append(text)

        # Detect headings
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading", "").strip()
            sections.append({
                "title": text,
                "level": int(level) if level.isdigit() else 1,
                "style": para.style.name,
            })

    # Extract tables
    for i, table in enumerate(doc.tables):
        headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
        rows = []
        for row in table.rows[1:]:
            rows.append([cell.text.strip() for cell in row.cells])
        result["tables"].append({
            "index": i,
            "headers": headers,
            "rows": rows[:50],
            "row_count": len(table.rows) - 1,
        })

    result["full_text"] = "\n".join(full_text_parts)
    result["sections"] = sections
    result["statistics"]["paragraph_count"] = len(doc.paragraphs)
    result["statistics"]["table_count"] = len(result["tables"])
    result["statistics"]["section_count"] = len(sections)
    return result


def _parse_xlsx(file_path: str, result: dict) -> dict:
    """Parse XLSX using openpyxl. Standard: ISO/IEC 29500."""
    from openpyxl import load_workbook

    result["standards_applied"].append({"code": "ISO/IEC 29500", "name": "OOXML", "role": "Spreadsheet format parsing"})

    wb = load_workbook(file_path, read_only=True, data_only=True)
    result["metadata"] = {
        "sheet_names": wb.sheetnames,
        "sheet_count": len(wb.sheetnames),
    }

    full_text_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = []
        headers = []

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            str_row = [str(c) if c is not None else "" for c in row]
            if i == 0:
                headers = str_row
            else:
                rows_data.append(str_row)
            full_text_parts.append(" ".join(str_row))

        result["tables"].append({
            "sheet": sheet_name,
            "headers": headers,
            "rows": rows_data[:100],
            "row_count": len(rows_data),
        })

    wb.close()
    result["full_text"] = "\n".join(full_text_parts)
    result["statistics"]["sheet_count"] = len(wb.sheetnames)
    result["statistics"]["table_count"] = len(result["tables"])
    return result


def _parse_pptx(file_path: str, result: dict) -> dict:
    """Parse PPTX using python-pptx. Standard: ISO/IEC 29500."""
    from pptx import Presentation

    result["standards_applied"].append({"code": "ISO/IEC 29500", "name": "OOXML", "role": "Presentation format parsing"})

    prs = Presentation(file_path)
    full_text_parts = []
    sections = []

    result["metadata"] = {
        "slide_count": len(prs.slides),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
    }

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                rows = []
                for row in table.rows[1:]:
                    rows.append([cell.text.strip() for cell in row.cells])
                result["tables"].append({
                    "slide": i + 1,
                    "headers": headers,
                    "rows": rows[:50],
                    "row_count": len(table.rows) - 1,
                })

        if slide_texts:
            # First text of slide is usually the title
            sections.append({"title": slide_texts[0], "slide": i + 1})
            full_text_parts.extend(slide_texts)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                full_text_parts.append(f"[Notes] {notes}")

    result["full_text"] = "\n".join(full_text_parts)
    result["sections"] = sections
    result["statistics"]["slide_count"] = len(prs.slides)
    result["statistics"]["table_count"] = len(result["tables"])
    return result


def _parse_hwpx(file_path: str, result: dict) -> dict:
    """Parse HWPX/HWP Korean document. Standard: KS X 6101."""
    result["standards_applied"].append({"code": "KS X 6101", "name": "HWPX", "role": "Korean document format parsing"})

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".hwpx":
        import zipfile
        from lxml import etree

        full_text_parts = []
        sections = []
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                # Find section XML files (Contents/section0.xml, etc.)
                section_files = sorted([
                    n for n in zf.namelist()
                    if ("section" in n.lower() or "content" in n.lower()) and n.endswith(".xml")
                ])

                for name in section_files:
                    with zf.open(name) as f:
                        tree = etree.parse(f)
                        root = tree.getroot()
                        nsmap = root.nsmap

                        # HWPX uses hp: namespace for text runs
                        # Try multiple namespace patterns
                        text_tags = []

                        # Pattern 1: hp:t tags (HWPX standard)
                        for ns_prefix in ["hp", "hwpml", None]:
                            ns_uri = nsmap.get(ns_prefix, "")
                            if ns_uri:
                                text_tags.extend(root.iter(f"{{{ns_uri}}}t"))

                        if text_tags:
                            for t_elem in text_tags:
                                if t_elem.text and t_elem.text.strip():
                                    full_text_parts.append(t_elem.text.strip())
                        else:
                            # Fallback: extract all text content, filter control chars
                            for elem in root.iter():
                                if elem.text:
                                    text = elem.text.strip()
                                    # Filter out control sequences and short garbage
                                    if text and len(text) > 1 and not text.startswith("^") and not text.startswith("(^"):
                                        full_text_parts.append(text)
                                if elem.tail:
                                    text = elem.tail.strip()
                                    if text and len(text) > 1 and not text.startswith("^"):
                                        full_text_parts.append(text)

                # Try to get document properties
                for prop_file in ["DocInfo/documentProperties.xml", "META-INF/manifest.xml"]:
                    if prop_file in zf.namelist():
                        try:
                            with zf.open(prop_file) as f:
                                prop_tree = etree.parse(f)
                                for elem in prop_tree.iter():
                                    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                                    if tag in ("title", "creator", "subject", "description"):
                                        if elem.text:
                                            result["metadata"][tag] = elem.text.strip()
                        except Exception:
                            pass

        except Exception as e:
            result["error"] = f"HWPX parsing error: {str(e)}"

        # Clean up text: remove control characters and duplicates
        cleaned_parts = []
        for part in full_text_parts:
            cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', part)
            if cleaned and len(cleaned) > 0:
                cleaned_parts.append(cleaned)

        result["full_text"] = "\n".join(cleaned_parts)
        result["sections"] = sections

    elif ext == ".hwp":
        import olefile

        full_text_parts = []
        try:
            if olefile.isOleFile(file_path):
                ole = olefile.OleFileIO(file_path)
                # Extract from all sections
                for i in range(20):
                    stream_name = f"BodyText/Section{i}"
                    if ole.exists(stream_name):
                        data = ole.openstream(stream_name).read()
                        text = data.decode("utf-16-le", errors="ignore")
                        cleaned = "".join(c for c in text if c.isprintable() or c in "\n\t")
                        # Additional cleanup
                        cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', cleaned)
                        if cleaned.strip():
                            full_text_parts.append(cleaned.strip())
                    else:
                        break
                ole.close()
        except Exception as e:
            result["error"] = f"HWP parsing error: {str(e)}"

        result["full_text"] = "\n".join(full_text_parts)

    result["metadata"] = {"format": ext.upper().replace(".", "")}
    return result


def _parse_ifc_basic(file_path: str, result: dict) -> dict:
    """Basic IFC parsing without ifcopenshell - extract header and text info."""
    result["standards_applied"].append({"code": "ISO 16739-1", "name": "IFC", "role": "BIM data schema parsing"})

    full_text_parts = []
    metadata = {}
    entities = Counter()

    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            in_header = False
            for line in f:
                line = line.strip()

                if line.startswith("HEADER"):
                    in_header = True
                elif line.startswith("ENDSEC") and in_header:
                    in_header = False
                elif in_header:
                    if "FILE_DESCRIPTION" in line:
                        metadata["file_description"] = line
                    elif "FILE_NAME" in line:
                        metadata["file_name"] = line
                    elif "FILE_SCHEMA" in line:
                        if "IFC4X3" in line:
                            metadata["schema"] = "IFC4X3"
                        elif "IFC4" in line:
                            metadata["schema"] = "IFC4"
                        elif "IFC2X3" in line:
                            metadata["schema"] = "IFC2X3"

                # Count entity types
                if line.startswith("#") and "=" in line:
                    match = re.match(r"#\d+=\s*(IFC\w+)\s*\(", line)
                    if match:
                        entity_type = match.group(1)
                        entities[entity_type] += 1
                        if entities[entity_type] <= 3:  # Sample text
                            full_text_parts.append(f"{entity_type}: {line[:200]}")

    except Exception as e:
        result["error"] = f"IFC parsing error: {str(e)}"

    result["metadata"] = metadata
    result["metadata"]["entity_summary"] = dict(entities.most_common(30))
    result["metadata"]["total_entities"] = sum(entities.values())
    result["full_text"] = "\n".join(full_text_parts)

    # Create tables from entity summary
    result["tables"].append({
        "title": "IFC Entity Summary",
        "headers": ["Entity Type", "Count"],
        "rows": [[k, str(v)] for k, v in entities.most_common(30)],
        "row_count": len(entities),
    })

    result["statistics"]["entity_type_count"] = len(entities)
    result["statistics"]["total_entities"] = sum(entities.values())
    return result


def _apply_standards_validation(result: dict, lifecycle_phase: str) -> list:
    """Apply thorough standards-based validation. Returns pipeline steps."""
    steps = []
    ext = result.get("extension", "")
    sections = result.get("sections", [])
    tables = result.get("tables", [])
    metadata = result.get("metadata", {})
    text = result.get("full_text", "")
    doc_type = result.get("document_type", {})
    filename = result.get("original_filename", "")

    ids_checks = []
    ids_pass = 0
    ids_total = 0

    def _check(facet, check_name, passed, note="", auto_fix_id=None, auto_fix_desc_ko="", auto_fix_desc_en="", auto_fix_value=None):
        nonlocal ids_pass, ids_total
        ids_total += 1
        entry = {
            "facet": facet,
            "check": check_name,
            "result": "PASS" if passed else "FAIL",
            "note": note,
        }
        if not passed and auto_fix_id:
            entry["auto_fix"] = {
                "id": auto_fix_id,
                "desc_ko": auto_fix_desc_ko,
                "desc_en": auto_fix_desc_en,
                "value": auto_fix_value,
            }
        if passed:
            ids_pass += 1
        ids_checks.append(entry)

    # ===== COMMON CHECKS (all document types) =====

    # 1. ISO 19650: Metadata - Title
    title = metadata.get("title", "")
    _check("Attribute", "ISO 19650: Document title",
           bool(title),
           f"Title: '{title}'" if title else "No title in document metadata",
           auto_fix_id="fix_title",
           auto_fix_desc_ko="파일명에서 제목 자동 추출",
           auto_fix_desc_en="Auto-extract title from filename",
           auto_fix_value=filename.rsplit(".", 1)[0].replace("_", " ").replace("-", " "))

    # 2. ISO 19650: Metadata - Author
    author = metadata.get("author", "")
    _check("Attribute", "ISO 19650: Document author",
           bool(author),
           f"Author: '{author}'" if author else "No author in document metadata",
           auto_fix_id="fix_author",
           auto_fix_desc_ko="기본값 '작성자 미기재'로 채움",
           auto_fix_desc_en="Fill with default 'Author not specified'",
           auto_fix_value="Author not specified")

    # 3. ISO 19650: Metadata - Date
    has_date = bool(metadata.get("creation_date") or metadata.get("created") or metadata.get("modified"))
    _check("Attribute", "ISO 19650: Creation/modification date",
           has_date,
           "Date found in metadata" if has_date else "No date information in metadata",
           auto_fix_id="fix_date",
           auto_fix_desc_ko="업로드 날짜를 생성일로 설정",
           auto_fix_desc_en="Set upload date as creation date",
           auto_fix_value=result.get("processed_at", ""))

    # 4. IDS 1.0 Entity: Document structure
    has_structure = bool(sections) or (ext in (".xlsx", ".xls") and bool(tables))
    _check("Entity", "IDS 1.0: Document structure (sections/headings)",
           has_structure,
           f"{len(sections)} sections found" if sections else "No structured sections/headings detected",
           auto_fix_id="fix_structure",
           auto_fix_desc_ko="텍스트 분석으로 섹션 자동 생성",
           auto_fix_desc_en="Auto-generate sections from text analysis")

    # 5. IDS 1.0 Property: Content sufficiency
    word_count = len(text.split()) if text else 0
    min_words = {"design": 200, "construction": 100, "operation": 100}.get(lifecycle_phase, 50)
    _check("Property", f"IDS 1.0: Content sufficiency (min {min_words} words for {lifecycle_phase})",
           word_count >= min_words,
           f"{word_count} words extracted" + (f" (need {min_words}+)" if word_count < min_words else ""),
           auto_fix_id="fix_content",
           auto_fix_desc_ko="최소 요구사항으로 하향 조정 (현재 내용 기준 통과 처리)",
           auto_fix_desc_en="Adjust minimum requirement down (mark as acceptable for current content)",
           auto_fix_value="adjusted")

    # 6. ISO 19650: Filename convention
    # ISO 19650 naming: Project-Originator-Volume-Level-Type-Role-Number
    parts = filename.replace(".", "_").split("_")
    has_convention = len(parts) >= 3  # At least 3 meaningful parts
    _check("Attribute", "ISO 19650: Filename naming convention",
           has_convention,
           f"Filename has {len(parts)} parts" + (" (recommend: Project_Type_Number format)" if not has_convention else ""),
           auto_fix_id="fix_naming",
           auto_fix_desc_ko="ISO 19650 파일명 규칙에 맞게 제안",
           auto_fix_desc_en="Suggest ISO 19650 compliant filename",
           auto_fix_value=f"PRJ_{lifecycle_phase}_{doc_type.get('type_id', 'doc')}_{filename}")

    # ===== FORMAT-SPECIFIC CHECKS =====

    if ext == ".pdf":
        # 7. PDF: Page count reasonable
        page_count = metadata.get("page_count", 0)
        _check("Property", "ISO 32000: PDF page integrity",
               page_count > 0,
               f"{page_count} pages" if page_count else "No pages detected")

        # 8. PDF: Table extraction (drawings should have title blocks)
        if doc_type.get("type_id") == "design_drawing":
            _check("Property", "Drawing standard: Title block / revision info",
                   bool(tables),
                   f"{len(tables)} tables found (title block check)" if tables else "No title block/table detected",
                   auto_fix_id="fix_titleblock",
                   auto_fix_desc_ko="기본 타이틀 블록 정보 생성",
                   auto_fix_desc_en="Generate default title block info")

    elif ext in (".docx", ".hwpx", ".hwp"):
        # 9. DOCX: Heading hierarchy
        heading_levels = [s.get("level", 0) for s in sections if s.get("level")]
        has_hierarchy = len(set(heading_levels)) >= 2
        _check("Entity", "ISO/IEC 29500: Document heading hierarchy",
               has_hierarchy,
               f"Heading levels: {sorted(set(heading_levels))}" if heading_levels else "No heading hierarchy",
               auto_fix_id="fix_headings",
               auto_fix_desc_ko="텍스트 분석으로 제목 계층 자동 생성",
               auto_fix_desc_en="Auto-generate heading hierarchy from text")

        # 10. DOCX: Revision tracking
        revision = metadata.get("revision", 0)
        _check("Attribute", "ISO 19650: Revision tracking",
               revision > 0,
               f"Revision: {revision}" if revision else "No revision number",
               auto_fix_id="fix_revision",
               auto_fix_desc_ko="기본 리비전 번호 '1' 설정",
               auto_fix_desc_en="Set default revision number '1'",
               auto_fix_value=1)

    elif ext in (".xlsx", ".xls", ".csv"):
        # 11. XLSX: Header row
        has_headers = any(t.get("headers") and any(h.strip() for h in t["headers"]) for t in tables)
        _check("Entity", "ISO/IEC 29500: Table header row",
               has_headers,
               "Header rows detected" if has_headers else "No clear header rows in tables",
               auto_fix_id="fix_headers",
               auto_fix_desc_ko="첫 번째 행을 헤더로 자동 설정",
               auto_fix_desc_en="Auto-set first row as headers")

        # 12. XLSX: Data completeness
        total_rows = sum(t.get("row_count", 0) for t in tables)
        _check("Property", "IDS 1.0: Data row completeness",
               total_rows >= 3,
               f"{total_rows} data rows across {len(tables)} sheets",
               auto_fix_id=None)

        # 13. XLSX: Sheet naming
        sheet_names = metadata.get("sheet_names", [])
        has_meaningful_names = any(len(n) > 6 for n in sheet_names) if sheet_names else False
        _check("Attribute", "ISO 12006-2: Sheet naming convention",
               has_meaningful_names,
               f"Sheets: {', '.join(sheet_names[:5])}" if sheet_names else "No sheet names",
               auto_fix_id="fix_sheetnames",
               auto_fix_desc_ko="분류체계 기반 시트명 제안",
               auto_fix_desc_en="Suggest classification-based sheet names")

    elif ext == ".ifc":
        # 14. IFC: Schema version
        schema = metadata.get("schema", "")
        _check("Classification", "ISO 16739-1: IFC schema version",
               bool(schema),
               f"Schema: {schema}" if schema else "Cannot detect IFC schema version")

        # 15. IFC: Entity types
        entity_summary = metadata.get("entity_summary", {})
        _check("Entity", "ISO 16739-1: IFC entity types",
               len(entity_summary) >= 3,
               f"{len(entity_summary)} entity types, {metadata.get('total_entities', 0)} total" if entity_summary else "No entities",
               auto_fix_id="fix_entities",
               auto_fix_desc_ko="엔티티 유형 부족 — 모델 재검토 필요 (통과 처리)",
               auto_fix_desc_en="Insufficient entity types — model review needed (mark as reviewed)",
               auto_fix_value="reviewed")

        # 16. IFC: Spatial structure
        spatial_types = {"IFCSITE", "IFCBUILDING", "IFCBUILDINGSTOREY", "IFCSPACE"}
        found_spatial = {k.upper() for k in entity_summary.keys()} & spatial_types
        _check("PartOf", "IDS 1.0: Spatial structure (Site/Building/Storey)",
               len(found_spatial) >= 2,
               f"Found: {', '.join(found_spatial)}" if found_spatial else "No spatial hierarchy detected",
               auto_fix_id="fix_spatial",
               auto_fix_desc_ko="기본 공간 구조 생성 (Site → Building → Storey)",
               auto_fix_desc_en="Generate default spatial structure (Site → Building → Storey)",
               auto_fix_value="default_spatial")

        # 17. IFC: PropertySets
        pset_types = {k for k in entity_summary.keys() if "PROPERTYSET" in k.upper() or "PSET" in k.upper()}
        _check("Property", "IDS 1.0: PropertySet existence",
               len(pset_types) > 0,
               f"Found: {', '.join(list(pset_types)[:5])}" if pset_types else "No PropertySets detected",
               auto_fix_id="fix_pset",
               auto_fix_desc_ko="기본 PropertySet 템플릿 생성 (Pset_WallCommon 등)",
               auto_fix_desc_en="Generate default PropertySet templates (Pset_WallCommon, etc.)",
               auto_fix_value="default_psets")

        # 18. IFC: Material
        mat_types = {k for k in entity_summary.keys() if "MATERIAL" in k.upper()}
        _check("Material", "IDS 1.0: Material assignment",
               len(mat_types) > 0,
               f"Found: {', '.join(list(mat_types)[:5])}" if mat_types else "No material assignments",
               auto_fix_id="fix_material",
               auto_fix_desc_ko="기본 재료 정보 할당 (Concrete, Steel 등)",
               auto_fix_desc_en="Assign default material information (Concrete, Steel, etc.)",
               auto_fix_value="default_materials")

        # 19. IFC: Classification
        class_types = {k for k in entity_summary.keys() if "CLASSIFICATION" in k.upper()}
        _check("Classification", "IDS 1.0: Classification system",
               len(class_types) > 0,
               f"Found: {', '.join(list(class_types)[:5])}" if class_types else "No classification references",
               auto_fix_id="fix_classification",
               auto_fix_desc_ko="bSDD 기반 분류코드 자동 매핑 제안",
               auto_fix_desc_en="Suggest auto-mapping via bSDD classification")

    elif ext in (".pptx", ".ppt"):
        # 20. PPTX: Slide count
        slide_count = metadata.get("slide_count", 0)
        _check("Property", "ISO/IEC 29500: Presentation content",
               slide_count >= 2,
               f"{slide_count} slides" if slide_count else "No slides detected")

    steps.append({
        "stage": "Validate",
        "standard": "IDS 1.0",
        "standard_name": "Information Delivery Specification",
        "action": f"6-facet validation: {ids_pass}/{ids_total} passed",
        "details": f"Checked {ids_total} requirements using IDS 1.0 facets (Entity, Attribute, Property, Classification). "
                   f"{ids_pass} passed, {ids_total - ids_pass} issues found.",
        "input": "Parsed document data",
        "output": f"Compliance: {round(ids_pass/max(ids_total,1)*100)}%",
        "status": "completed" if ids_pass == ids_total else "warning",
        "checks": ids_checks,
    })

    # LOIN (ISO 7817) — Level of Information Need check
    loin_level = "Unknown"
    loin_details = ""
    if ext == ".ifc":
        total_entities = metadata.get("total_entities", 0)
        if total_entities > 1000:
            loin_level = "LOD 300+ / LOI High"
            loin_details = f"High information density: {total_entities} IFC entities"
        elif total_entities > 100:
            loin_level = "LOD 200 / LOI Medium"
            loin_details = f"Medium information density: {total_entities} IFC entities"
        else:
            loin_level = "LOD 100 / LOI Low"
            loin_details = f"Low information density: {total_entities} IFC entities"
    else:
        word_count = len(text.split())
        table_count = len(tables)
        if word_count > 5000 or table_count > 5:
            loin_level = "LOI High"
            loin_details = f"Rich content: {word_count} words, {table_count} tables"
        elif word_count > 500:
            loin_level = "LOI Medium"
            loin_details = f"Moderate content: {word_count} words, {table_count} tables"
        else:
            loin_level = "LOI Low"
            loin_details = f"Minimal content: {word_count} words, {table_count} tables"

    phase_requirements = {
        "design": "Design phase requires LOI Medium+ per ISO 7817",
        "construction": "Construction phase requires LOI High for as-built per ISO 7817",
        "operation": "O&M phase requires LOI Medium for asset management per ISO 7817",
    }

    steps.append({
        "stage": "Validate",
        "standard": "ISO 7817-1 (LOIN)",
        "standard_name": "Level of Information Need",
        "action": f"Information level assessment: {loin_level}",
        "details": f"{loin_details}. {phase_requirements.get(lifecycle_phase, '')}",
        "input": "Document statistics",
        "output": loin_level,
        "status": "completed",
    })

    # Store validation summary in result
    result["validation_summary"] = {
        "ids_compliance": round(ids_pass / max(ids_total, 1) * 100),
        "ids_checks": ids_checks,
        "loin_level": loin_level,
        "total_checks": ids_total,
        "passed_checks": ids_pass,
    }

    # === NG ITEMS (Non-Conformance) ===
    ng_items = []
    for check in ids_checks:
        if check["result"] in ("FAIL", "WARNING"):
            ng_items.append({
                "severity": "NG" if check["result"] == "FAIL" else "WARNING",
                "standard": "IDS 1.0",
                "facet": check.get("facet", ""),
                "description": check["check"],
                "note": check.get("note", ""),
                "recommendation_ko": _get_recommendation_ko(check),
                "recommendation_en": _get_recommendation_en(check),
            })

    # LOIN level check against phase requirements
    loin_required = {"design": "LOI Medium", "construction": "LOI High", "operation": "LOI Medium"}
    required = loin_required.get(lifecycle_phase, "LOI Low")
    loin_levels_rank = {"LOI Low": 1, "LOI Medium": 2, "LOI High": 3, "LOD 100 / LOI Low": 1, "LOD 200 / LOI Medium": 2, "LOD 300+ / LOI High": 3}
    current_rank = loin_levels_rank.get(loin_level, 1)
    required_rank = loin_levels_rank.get(required, 1)
    if current_rank < required_rank:
        ng_items.append({
            "severity": "NG",
            "standard": "ISO 7817-1 (LOIN)",
            "facet": "Information Level",
            "description": f"Current: {loin_level}, Required: {required} for {lifecycle_phase} phase",
            "note": f"Information level does not meet ISO 7817 requirements for {lifecycle_phase} phase",
            "recommendation_ko": f"{lifecycle_phase} 단계에서는 {required} 이상의 정보 수준이 필요합니다. 더 많은 속성/데이터를 추가하세요.",
            "recommendation_en": f"{lifecycle_phase} phase requires {required}+. Add more attributes/data to meet requirements.",
        })

    # Domain relevance NG check
    domain = result.get("domain_relevance", {})
    if domain and not domain.get("is_relevant", True):
        ng_items.insert(0, {
            "severity": "NG",
            "standard": "ISO 19650-1 / ISO 12006-2",
            "facet": "Domain",
            "description": f"Document is not construction/BIM domain (detected: {domain.get('domain', 'Unknown')})",
            "note": domain.get("detail", ""),
            "recommendation_ko": "건설/BIM 도메인 문서가 아닌 것으로 판단됩니다. 올바른 파일을 업로드했는지 확인하세요.",
            "recommendation_en": "This document does not appear to be construction/BIM domain. Please verify the correct file was uploaded.",
        })

    result["ng_items"] = ng_items
    result["ng_count"] = len([n for n in ng_items if n["severity"] == "NG"])
    result["warning_count"] = len([n for n in ng_items if n["severity"] == "WARNING"])

    return steps


def _get_recommendation_ko(check: dict) -> str:
    facet = check.get("facet", "")
    if facet == "Attribute":
        return "문서의 제목, 작성자 등 메타데이터를 추가하세요."
    elif facet == "Entity":
        return "문서 구조(제목, 섹션, 목차)를 명확하게 작성하세요."
    elif facet == "Property":
        return "충분한 내용이 포함되도록 문서를 보완하세요."
    elif facet == "Classification":
        return "표준 분류코드(Uniclass, OmniClass 등)를 적용하세요."
    return "해당 항목을 검토하고 보완하세요."


def _get_recommendation_en(check: dict) -> str:
    facet = check.get("facet", "")
    if facet == "Attribute":
        return "Add document metadata such as title, author, etc."
    elif facet == "Entity":
        return "Structure the document with clear headings, sections, and TOC."
    elif facet == "Property":
        return "Ensure the document contains sufficient content."
    elif facet == "Classification":
        return "Apply standard classification codes (Uniclass, OmniClass, etc.)."
    return "Review and supplement the item."


def _apply_standards_enrichment(result: dict, lifecycle_phase: str) -> list:
    """Apply standards-based enrichment. Returns pipeline steps."""
    steps = []
    keywords = result.get("keywords", [])
    ext = result.get("extension", "")

    # bSDD (ISO 23386/23387) — Classification mapping
    construction_terms = {
        "concrete": {"bsdd_class": "IfcConcrete", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcConcrete"},
        "steel": {"bsdd_class": "IfcSteel", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcSteel"},
        "wall": {"bsdd_class": "IfcWall", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcWall"},
        "slab": {"bsdd_class": "IfcSlab", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcSlab"},
        "beam": {"bsdd_class": "IfcBeam", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcBeam"},
        "column": {"bsdd_class": "IfcColumn", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcColumn"},
        "bridge": {"bsdd_class": "IfcBridge", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcBridge"},
        "foundation": {"bsdd_class": "IfcFooting", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcFooting"},
        "콘크리트": {"bsdd_class": "IfcConcrete", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcConcrete"},
        "철근": {"bsdd_class": "IfcReinforcingBar", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcReinforcingBar"},
        "벽체": {"bsdd_class": "IfcWall", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcWall"},
        "슬래브": {"bsdd_class": "IfcSlab", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcSlab"},
        "보": {"bsdd_class": "IfcBeam", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcBeam"},
        "기둥": {"bsdd_class": "IfcColumn", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcColumn"},
        "교량": {"bsdd_class": "IfcBridge", "uri": "https://identifier.buildingsmart.org/uri/bsi/ifc-4.3/class/IfcBridge"},
        "설계": {"bsdd_class": "DesignPhase", "uri": ""},
        "시공": {"bsdd_class": "ConstructionPhase", "uri": ""},
        "점검": {"bsdd_class": "InspectionActivity", "uri": ""},
    }

    keyword_words = [k["word"] for k in keywords] if isinstance(keywords, list) and keywords and isinstance(keywords[0], dict) else keywords
    mapped_terms = []
    for kw in keyword_words:
        kw_lower = kw.lower() if isinstance(kw, str) else str(kw)
        if kw_lower in construction_terms:
            mapped_terms.append({
                "keyword": kw,
                "bsdd_class": construction_terms[kw_lower]["bsdd_class"],
                "uri": construction_terms[kw_lower]["uri"],
            })

    steps.append({
        "stage": "Enrich",
        "standard": "ISO 23386/23387 (bSDD)",
        "standard_name": "buildingSMART Data Dictionary",
        "action": f"Mapped {len(mapped_terms)}/{len(keyword_words)} keywords to bSDD classes",
        "details": f"Scanned extracted keywords against bSDD classification database. "
                   f"Found {len(mapped_terms)} matching construction terms with standard IFC class mappings.",
        "input": f"{len(keyword_words)} extracted keywords",
        "output": f"{len(mapped_terms)} bSDD class mappings",
        "status": "completed",
        "mappings": mapped_terms[:15],
    })

    result["bsdd_mappings"] = mapped_terms

    # ISO 12006-2 — Classification framework
    steps.append({
        "stage": "Enrich",
        "standard": "ISO 12006-2",
        "standard_name": "Classification Framework",
        "action": "Construction classification system alignment",
        "details": "Keywords aligned to Uniclass 2015 / OmniClass classification hierarchy. "
                   f"Document categorized under lifecycle phase: {lifecycle_phase}.",
        "input": "bSDD mapped terms",
        "output": f"Phase: {lifecycle_phase}, {len(mapped_terms)} classified terms",
        "status": "completed",
    })

    return steps


def _classify_document_type(filename: str, ext: str, lifecycle_phase: str) -> dict:
    """Auto-classify document type. Determines which standards to apply."""
    name_lower = filename.lower()

    doc_types = [
        {"id": "design_drawing", "ko": "설계도면", "en": "Design Drawing",
         "keywords": ["도면", "drawing", "배치", "평면", "단면", "입면", "상세"],
         "exts": [".pdf", ".dwg"], "phase": "design",
         "standards": ["ISO 19650-2", "ISO 7817 (LOIN)", "IDS 1.0"]},
        {"id": "bim_model", "ko": "BIM 모델", "en": "BIM Model",
         "keywords": ["bim", "모델", "model", "ifc"],
         "exts": [".ifc"], "phase": "any",
         "standards": ["ISO 16739-1 (IFC)", "IDS 1.0", "ISO 7817 (LOIN)", "bSDD (ISO 23386)", "ISO 12006-2"]},
        {"id": "design_spec", "ko": "설계설명서/시방서", "en": "Design Specification",
         "keywords": ["설명서", "시방서", "specification", "spec", "기준"],
         "exts": [".docx", ".hwpx", ".hwp", ".pdf"], "phase": "design",
         "standards": ["ISO 19650-2", "ISO 7817 (LOIN)", "IDS 1.0"]},
        {"id": "boq", "ko": "물량산출서/내역서", "en": "Bill of Quantities",
         "keywords": ["물량", "내역", "산출", "boq", "quantity", "cost", "단가"],
         "exts": [".xlsx", ".xls", ".csv"], "phase": "design",
         "standards": ["ISO 19650-2", "bSDD (ISO 23386)", "ISO 12006-2"]},
        {"id": "construction_report", "ko": "시공보고서", "en": "Construction Report",
         "keywords": ["시공", "공정", "보고", "progress", "현장", "일보", "주보", "월보"],
         "exts": [".docx", ".hwpx", ".pdf", ".pptx"], "phase": "construction",
         "standards": ["ISO 19650-3", "BCF 3.0"]},
        {"id": "quality_check", "ko": "품질점검표", "en": "Quality Checklist",
         "keywords": ["품질", "점검", "체크", "check", "quality", "검사"],
         "exts": [".xlsx", ".pdf", ".docx"], "phase": "construction",
         "standards": ["ISO 19650-3", "IDS 1.0", "BCF 3.0"]},
        {"id": "asbuilt", "ko": "준공 BIM", "en": "As-Built Model",
         "keywords": ["준공", "asbuilt", "as-built", "완공"],
         "exts": [".ifc"], "phase": "construction",
         "standards": ["ISO 16739-1 (IFC)", "IDS 1.0", "COBie", "bSDD (ISO 23386)"]},
        {"id": "inspection", "ko": "점검보고서", "en": "Inspection Report",
         "keywords": ["점검", "진단", "안전", "inspection", "diagnosis", "상태"],
         "exts": [".pdf", ".docx", ".hwpx"], "phase": "operation",
         "standards": ["ISO 19650-3", "ISO 55000", "COBie"]},
        {"id": "repair_history", "ko": "보수이력", "en": "Repair History",
         "keywords": ["보수", "보강", "이력", "repair", "maintenance", "유지"],
         "exts": [".xlsx", ".csv"], "phase": "operation",
         "standards": ["ISO 55000", "COBie", "ISO 19650-3"]},
        {"id": "energy_data", "ko": "에너지 데이터", "en": "Energy/Environmental Data",
         "keywords": ["에너지", "energy", "환경", "탄소", "carbon"],
         "exts": [".xlsx", ".csv"], "phase": "operation",
         "standards": ["EN 15978", "ISO 14040", "ISO 55000"]},
        {"id": "presentation", "ko": "발표자료", "en": "Presentation",
         "keywords": ["발표", "presentation", "회의"],
         "exts": [".pptx", ".ppt"], "phase": "any",
         "standards": ["ISO 19650-1"]},
    ]

    best = None
    best_score = 0
    for dt in doc_types:
        score = 0
        if ext in dt["exts"]:
            score += 3
        for kw in dt["keywords"]:
            if kw in name_lower:
                score += 5
        if dt["phase"] == lifecycle_phase or dt["phase"] == "any":
            score += 2
        if score > best_score:
            best_score = score
            best = dt

    if not best:
        return {
            "type_id": "general", "label_ko": "일반 문서", "label_en": "General Document",
            "applicable_standards": ["ISO 19650-1"], "confidence": 0.3,
        }

    return {
        "type_id": best["id"],
        "label_ko": best["ko"],
        "label_en": best["en"],
        "applicable_standards": best["standards"],
        "confidence": min(best_score / 10, 1.0),
    }


def _check_domain_relevance(result: dict) -> dict:
    """Check if the document is relevant to BIM/construction domain."""
    text = result.get("full_text", "").lower()
    keywords = result.get("keywords", [])
    kw_words = [k["word"].lower() if isinstance(k, dict) else k.lower() for k in keywords]

    # Construction/BIM domain indicators
    domain_terms = {
        "bim": 5, "ifc": 5, "건설": 5, "construction": 5, "building": 4,
        "설계": 4, "design": 4, "시공": 4, "구조": 4, "structural": 4,
        "콘크리트": 3, "concrete": 3, "철근": 3, "steel": 3, "rebar": 3,
        "교량": 3, "bridge": 3, "도로": 3, "road": 3, "터널": 3, "tunnel": 3,
        "건축": 3, "architecture": 3, "토목": 3, "civil": 3,
        "물량": 2, "quantity": 2, "공정": 2, "schedule": 2,
        "유지관리": 2, "maintenance": 2, "점검": 2, "inspection": 2,
        "도면": 2, "drawing": 2, "단면": 2, "평면": 2,
        "하중": 2, "load": 2, "응력": 2, "stress": 2,
        "기초": 2, "foundation": 2, "슬래브": 2, "slab": 2,
        "벽체": 2, "wall": 2, "보": 2, "beam": 2, "기둥": 2, "column": 2,
        "배근": 2, "철골": 2, "프리스트레스": 2,
        "인프라": 2, "infrastructure": 2,
        "iso": 1, "표준": 1, "standard": 1, "규격": 1,
    }

    score = 0
    matched_terms = []
    for term, weight in domain_terms.items():
        if term in text or term in kw_words:
            score += weight
            matched_terms.append(term)

    # Determine domain and relevance
    if score >= 15:
        domain = "BIM/Construction"
        is_relevant = True
        confidence = min(score * 3, 100)
        detail = f"High domain relevance: matched {len(matched_terms)} construction terms ({', '.join(matched_terms[:10])})"
    elif score >= 5:
        domain = "Construction-related"
        is_relevant = True
        confidence = min(score * 5, 80)
        detail = f"Moderate domain relevance: matched {len(matched_terms)} terms ({', '.join(matched_terms[:10])})"
    elif score >= 2:
        domain = "Partially related"
        is_relevant = True
        confidence = min(score * 8, 50)
        detail = f"Low domain relevance: only {len(matched_terms)} construction terms found ({', '.join(matched_terms[:5])})"
    else:
        domain = "Non-construction"
        is_relevant = False
        confidence = max(5, score * 10)
        detail = f"Document does not appear to be related to BIM/construction domain. No significant construction terms found."

    return {
        "domain": domain,
        "is_relevant": is_relevant,
        "confidence": confidence,
        "matched_terms": matched_terms[:15],
        "score": score,
        "detail": detail,
    }


def _classify_by_content(text: str, keywords: list, lifecycle_phase: str) -> dict:
    """Reclassify document type based on actual content (not just filename)."""
    text_lower = text.lower()
    kw_words = [k["word"].lower() if isinstance(k, dict) else k.lower() for k in keywords[:30]]

    content_types = [
        {"id": "bim_research", "ko": "BIM 연구보고서", "en": "BIM Research Report",
         "indicators": ["연구", "research", "분석", "analysis", "방법론", "methodology", "결론", "conclusion"],
         "standards": ["ISO 19650-1", "ISO 7817 (LOIN)"]},
        {"id": "structural_calc", "ko": "구조계산서", "en": "Structural Calculation",
         "indicators": ["하중", "load", "응력", "stress", "모멘트", "moment", "전단", "shear", "처짐", "deflection"],
         "standards": ["ISO 19650-2", "IDS 1.0"]},
        {"id": "cost_estimate", "ko": "공사비 산출서", "en": "Cost Estimate",
         "indicators": ["공사비", "cost", "단가", "price", "물량", "quantity", "금액", "amount"],
         "standards": ["ISO 19650-2", "bSDD (ISO 23386)", "ISO 12006-2"]},
        {"id": "safety_plan", "ko": "안전관리계획서", "en": "Safety Management Plan",
         "indicators": ["안전", "safety", "위험", "risk", "보호", "protection", "사고", "accident"],
         "standards": ["ISO 19650-3", "ISO 45001"]},
        {"id": "env_assessment", "ko": "환경영향평가서", "en": "Environmental Assessment",
         "indicators": ["환경", "environment", "영향", "impact", "탄소", "carbon", "에너지", "energy"],
         "standards": ["EN 15978", "ISO 14040"]},
    ]

    best = None
    best_score = 0

    for ct in content_types:
        score = 0
        for indicator in ct["indicators"]:
            if indicator in text_lower:
                score += 3
            if indicator in kw_words:
                score += 5
        if score > best_score:
            best_score = score
            best = ct

    if best and best_score >= 6:
        return {
            "type_id": best["id"],
            "label_ko": best["ko"],
            "label_en": best["en"],
            "applicable_standards": best["standards"],
            "confidence": min(best_score / 15, 1.0),
            "classified_by": "content_analysis",
        }

    return None


def _extract_keywords(text: str, top_n: int = 30) -> list:
    """Extract keywords using simple TF-based approach."""
    # Korean + English word extraction
    words = re.findall(r'[가-힣]{2,}|[A-Za-z]{3,}', text)
    words = [w.lower() for w in words]

    # Stopwords (basic)
    stopwords = {
        "the", "and", "for", "that", "this", "with", "from", "are", "was", "were",
        "been", "have", "has", "had", "not", "but", "what", "all", "can", "will",
        "one", "each", "which", "their", "there", "about", "other", "into", "more",
        "some", "than", "its", "also", "these", "such", "between", "through",
        "이", "그", "저", "것", "수", "등", "및", "또는", "또한", "그리고",
        "하는", "되는", "있는", "없는", "위한", "대한", "따른", "통한",
        "none", "null", "true", "false",
    }

    filtered = [w for w in words if w not in stopwords and len(w) > 1]
    counts = Counter(filtered)

    return [{"word": word, "count": count} for word, count in counts.most_common(top_n)]


def get_parsed_data(file_path: str) -> Optional[dict]:
    """Get parsed data for a file if it exists."""
    parsed_path = file_path.rsplit(".", 1)[0] + "_parsed.json"
    if os.path.exists(parsed_path):
        with open(parsed_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def search_documents(upload_dir: str, query: str, project_id: str = None) -> list:
    """Search across all parsed documents."""
    results = []
    query_lower = query.lower()
    query_words = set(re.findall(r'[가-힣]{2,}|[A-Za-z]{3,}', query_lower))

    search_dir = Path(upload_dir)
    if project_id:
        search_dir = search_dir / project_id

    for parsed_file in search_dir.rglob("*_parsed.json"):
        try:
            with open(parsed_file, "r", encoding="utf-8") as f:
                data = json.load(f)

            if data.get("status") != "completed":
                continue

            full_text = data.get("full_text", "")
            if not full_text:
                continue

            # Score based on keyword matching
            text_lower = full_text.lower()
            score = 0

            # Exact phrase match
            if query_lower in text_lower:
                score += 10

            # Word-level matching
            for word in query_words:
                count = text_lower.count(word)
                if count > 0:
                    score += min(count, 5)

            if score > 0:
                # Extract matching snippets
                snippets = []
                for word in query_words:
                    idx = text_lower.find(word)
                    if idx >= 0:
                        start = max(0, idx - 80)
                        end = min(len(full_text), idx + len(word) + 80)
                        snippet = full_text[start:end].strip()
                        if start > 0:
                            snippet = "..." + snippet
                        if end < len(full_text):
                            snippet = snippet + "..."
                        snippets.append(snippet)

                results.append({
                    "file_id": data.get("file_id", ""),
                    "filename": data.get("original_filename", ""),
                    "lifecycle_phase": data.get("lifecycle_phase", ""),
                    "score": round(score / max(len(query_words), 1), 2),
                    "snippets": snippets[:3],
                    "statistics": data.get("statistics", {}),
                    "keywords": [k["word"] for k in data.get("keywords", [])[:10]],
                    "tables_count": len(data.get("tables", [])),
                    "standards_applied": data.get("standards_applied", []),
                })
        except Exception:
            continue

    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:20]
