"""
PPTX Document Parser

Extracts text, tables, images, and metadata from Microsoft PowerPoint presentations.
Supports ISO/IEC 29500 (OOXML) format.
"""

import time
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime
import uuid

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .base import (
    DocumentParser,
    DocumentType,
    DocumentMetadata,
    ParseResult,
    Section,
    TableData,
)

logger = logging.getLogger(__name__)


class PPTXParser(DocumentParser):
    """
    Microsoft PowerPoint Presentation Parser

    Features:
    - Slide content extraction
    - Text from shapes, text boxes, tables
    - Speaker notes extraction
    - Image detection
    - Master slide and layout info
    - Embedded object detection
    """

    document_type = DocumentType.PPTX
    supported_extensions = ['.pptx', '.ppt']

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        self.extract_notes = self.config.get('extract_notes', True)
        self.extract_images = self.config.get('extract_images', True)

    async def parse(self, file_path: Path) -> ParseResult:
        """Parse PPTX and extract all content"""
        start_time = time.time()

        result = ParseResult(
            success=False,
            document_type=self.document_type,
            file_path=str(file_path),
            file_hash=self.compute_file_hash(file_path),
        )

        if not HAS_PPTX:
            result.errors.append("python-pptx not installed. Run: pip install python-pptx")
            return result

        try:
            prs = Presentation(file_path)

            # Extract metadata
            result.metadata = await self.extract_metadata(file_path)
            result.metadata.page_count = len(prs.slides)

            # Process slides
            sections = []
            all_tables = []
            text_parts = []
            table_idx = 0

            for slide_idx, slide in enumerate(prs.slides):
                slide_result = await self._process_slide(
                    slide, slide_idx, table_idx
                )

                # Create section from slide
                section = Section(
                    id=str(uuid.uuid4()),
                    heading=slide_result['title'] or f"Slide {slide_idx + 1}",
                    level=1,
                    content=slide_result['content'],
                    paragraphs=slide_result['paragraphs'],
                    tables=slide_result['tables'],
                )

                # Add speaker notes as subsection
                if slide_result['notes'] and self.extract_notes:
                    notes_section = Section(
                        id=str(uuid.uuid4()),
                        heading="Speaker Notes",
                        level=2,
                        content=slide_result['notes'],
                    )
                    section.children.append(notes_section)

                sections.append(section)
                all_tables.extend(slide_result['tables'])
                table_idx += len(slide_result['tables'])

                # Add to full text
                text_parts.append(f"=== Slide {slide_idx + 1}: {slide_result['title'] or 'Untitled'} ===")
                text_parts.append(slide_result['content'])
                if slide_result['notes']:
                    text_parts.append(f"[Notes: {slide_result['notes']}]")

            result.sections = sections
            result.tables = all_tables
            result.full_text = '\n\n'.join(text_parts)

            # Statistics
            word_count = len(result.full_text.split())
            result.statistics = {
                'slide_count': len(prs.slides),
                'table_count': len(all_tables),
                'character_count': len(result.full_text),
                'word_count': word_count,
            }
            result.metadata.word_count = word_count

            result.success = True

        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            result.errors.append(str(e))

        result.processing_time_ms = int((time.time() - start_time) * 1000)
        return result

    async def extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """Extract PPTX metadata"""
        metadata = DocumentMetadata()

        if not HAS_PPTX:
            return metadata

        try:
            prs = Presentation(file_path)
            props = prs.core_properties

            metadata.title = props.title
            metadata.author = props.author
            metadata.creator = props.author
            metadata.subject = props.subject

            if props.keywords:
                metadata.keywords = [k.strip() for k in props.keywords.split(',')]

            if props.created:
                metadata.created = props.created
            if props.modified:
                metadata.modified = props.modified

            metadata.page_count = len(prs.slides)

            metadata.custom = {
                'category': props.category,
                'comments': props.comments,
                'content_status': props.content_status,
                'last_modified_by': props.last_modified_by,
                'revision': props.revision,
                'version': props.version,
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height,
            }

        except Exception as e:
            logger.error(f"Error extracting PPTX metadata: {e}")

        return metadata

    async def extract_text(self, file_path: Path) -> str:
        """Extract full text from PPTX"""
        if not HAS_PPTX:
            return ""

        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide_idx, slide in enumerate(prs.slides):
                text_parts.append(f"=== Slide {slide_idx + 1} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

                # Speaker notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text:
                        text_parts.append(f"[Notes: {notes_text}]")

            return '\n'.join(text_parts)

        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""

    async def _process_slide(
        self,
        slide: 'pptx.slide.Slide',
        slide_idx: int,
        table_start_idx: int
    ) -> Dict[str, Any]:
        """Process a single slide"""
        result = {
            'title': None,
            'content': '',
            'paragraphs': [],
            'tables': [],
            'notes': '',
            'images': [],
            'shapes': [],
        }

        content_parts = []

        for shape in slide.shapes:
            # Title
            if shape.is_placeholder:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if hasattr(shape, 'placeholder_format'):
                    if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                        result['title'] = shape.text if hasattr(shape, 'text') else None

            # Text
            if hasattr(shape, "text") and shape.text:
                content_parts.append(shape.text)
                result['paragraphs'].append(shape.text)

            # Table
            if shape.has_table:
                table = shape.table
                table_data = TableData(
                    id=str(uuid.uuid4()),
                    name=f"Slide{slide_idx + 1}_Table{len(result['tables']) + 1}",
                    source_page=slide_idx + 1,
                )

                rows = []
                for row_idx, row in enumerate(table.rows):
                    row_data = [cell.text for cell in row.cells]
                    if row_idx == 0:
                        table_data.headers = row_data
                    else:
                        rows.append(row_data)

                table_data.rows = rows
                result['tables'].append(table_data)

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result['images'].append({
                    'name': shape.name,
                    'width': shape.width,
                    'height': shape.height,
                })

            # Track all shapes
            result['shapes'].append({
                'name': shape.name,
                'type': str(shape.shape_type),
                'has_text': hasattr(shape, 'text'),
            })

        result['content'] = '\n'.join(content_parts)

        # Speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            result['notes'] = notes_slide.notes_text_frame.text

        return result

    async def get_slide_thumbnails(
        self, file_path: Path
    ) -> List[Dict[str, Any]]:
        """Get information about slide layouts and thumbnails"""
        thumbnails = []

        if not HAS_PPTX:
            return thumbnails

        try:
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides):
                thumbnail_info = {
                    'slide_number': slide_idx + 1,
                    'layout_name': slide.slide_layout.name,
                    'shape_count': len(slide.shapes),
                    'has_notes': slide.has_notes_slide,
                }
                thumbnails.append(thumbnail_info)

        except Exception as e:
            logger.error(f"Error getting thumbnails: {e}")

        return thumbnails

    async def get_master_slides(self, file_path: Path) -> List[Dict[str, Any]]:
        """Get information about master slides"""
        masters = []

        if not HAS_PPTX:
            return masters

        try:
            prs = Presentation(file_path)

            for master in prs.slide_masters:
                master_info = {
                    'name': master.name if hasattr(master, 'name') else 'Unnamed',
                    'layouts': [
                        {
                            'name': layout.name,
                            'placeholder_count': len(layout.placeholders),
                        }
                        for layout in master.slide_layouts
                    ],
                }
                masters.append(master_info)

        except Exception as e:
            logger.error(f"Error getting master slides: {e}")

        return masters
